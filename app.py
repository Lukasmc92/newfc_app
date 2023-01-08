from inspect import stack
from bs4 import StopParsing
import streamlit as st
from calendar import month_abbr
import PIL
import lxml
import pandas as pd
import plotly
import plotly.express as px
from pyxlsb import open_workbook as open_xlsb
import plotly.graph_objs as go
import calendar
import numpy as np
import tempfile
from tempfile import NamedTemporaryFile
from plotly.graph_objects import Layout
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt
import io
from io import StringIO
from io import BytesIO
import codecs
import os
from pptx.dml.color import RGBColor
from contextlib import contextmanager
import dataframe_image as dfi
from subprocess import Popen
import plotly.io as pio
from kaleido.scopes.plotly import PlotlyScope
from streamlit_option_menu import option_menu
# from streamlit_lottie import st_lottie
import json
import xlrd



st.set_page_config(page_title="Food Cost Data App", 
        page_icon=":bar_chart:", 
        layout='wide')

pio.kaleido.scope.chromium_args = tuple([arg for arg in pio.kaleido.scope.chromium_args if arg != "--disable-dev-shm-usage"])


with st.sidebar:
    selected = option_menu(
        menu_title="Menu",
        options=["Main Page", "Food Cost %", "DataViz", "Instructions", "Contact"],
    )

if selected == "Main Page":
    st.image("./assets/Boarder.png")
    st.header("Welcome!")
    st.header("If you want to learn more about Food Cost in your restaurant, then you've come to the right place!")
    st.title("")
    st.title("")
    col1, col2, col3 = st.columns(3)
    with col1:
        st.subheader("Upload and analyze your Food Cost data:")
        st.image("./assets/Subcharts.png")
    with col1:
        st.image("./assets/bfstack.png")
    with col2:
        st.subheader("Calculate your Food Cost % throughout the month:")
        st.image("./assets/calc.png")
        st.image("./assets/stats.png")
    with col3:
        st.subheader('Download detailed reports to share with your team:')
        st.image("./assets/chickenwaste.png")
        st.image("./assets/Picture10.png")

if selected == "Instructions":
    
    st.title("Instructions for using this web application")
    st.markdown("---")
    st.subheader("What data should be uploaded?")
    st.write("From cfahome, donwload to your cmputer or phone the following:")
    st.subheader("ActivityReports (found in inform)")
    st.subheader("PMIX by date")
    st.subheader("Sales") 
    st.subheader("Sales and PMIX data must be downloaded from the data feed (search 'data feed' on cfahome)")
    
    st.write("Data feed files should be downloaded as csv's and should only include the date range you have completed inventory counts for.")
    st.write('Inventory counts are important because they allow you to create ActivityReports which supply the data needed to analyze food waste, and ultamtely, food cost percentage')
    st.subheader("What data should you be collecting?")
    st.write("Insight into food cost is best gained through data collection. The most important data you must collect is inventory data via the Inventory Activity Report. Month-to-month inventory activity reports can be downloaded from inForm.")
    st.subheader("How do I download the report?")
    st.write("The report is created when calculating the Food Cost %. A dwonload button will appear after you run the calculation")
    st.image('./assets/dlreport.png')
    st.subheader("Is it safe to upload my data?")
    st.write("Data uploaded to the site is not stored with an account of any kind. If data is uploaded, that data will dissapear as soon as you refresh or exit the page. Users accessing the site in at different sessions are not able to see what you upload.")
    st.write("This WebApp is built on top of Streamlit, a Python language framework. You can read more about Streamlit's security background in the links below")
    st.write("[Streamlit Cloud is now SOC 2 Type 1 compliant](https://blog.streamlit.io/streamlit-cloud-is-now-soc-2-type-1-compliant/)")
    st.write("[Streamlit Trust and Security](https://docs.streamlit.io/streamlit-cloud/trust-and-security)")
    st.write("[Where does st.file_uploader store uploaded files and when do they get deleted?](https://docs.streamlit.io/knowledge-base/using-streamlit/where-file-uploader-store-when-deleted)")
    col1, col2 = st.columns(2)
    st.markdown("---")
    st.subheader("Is there a better solution than Month-to-month data? Yes! Week-to-week inventory data! But wait...that means...")
    st.write("Yes, to keep up with your food cost, you should have weekly inventory counts (preferably after End-of-Day Saturday).")
    st.subheader("But isn't that labor intensive and time consuming?")
    st.write("This will take a little work, but not as much as you think. You don't need to count every single item like you do at end of month. To keep up with your food cost percentage, there are about 60 items I reccomend you count.")
    st.write("This counting process will take 30 minutes, depending on how well organized your inventory is.")

if selected == "Contact":
    st.title("Questions? Contact me!")
    st.subheader("Lukas McAnulty")
    st.subheader("Cell: (713) 992-7586")
    st.subheader("Email: lukas.mcanulty@gmail.com")

#buckets for catching dataframes and lists of dates
df_list = []
date_list = []
pmixdf = []
salesdf = []

st.cache()  
def sales_calc(dfsales):
            dfsales['Daily Sales'] = dfsales['Gross_Inside_Sales'] + dfsales['Gross_Outside_Sales'] - dfsales['Tax_Amount']
            dfsales['Date'] = pd.to_datetime(dfsales['Date'])
            dfsales['Date'] = dfsales['Date'].dt.date
            dfsales['Month'] = pd.DatetimeIndex(dfsales['Date']).month
            month = dfsales['Month Name'] = dfsales['Month'].apply(lambda x: calendar.month_name[x])
            salesdf.append(dfsales)
            global StopPgm
            StopPgm = ""
            return dfsales
st.cache()
def pmix_calc(dfpmix):
        dfpmix['Food Cost$'] = (dfpmix['Sold_Count']-dfpmix['Promotional_Count'])*dfpmix['Food_Cost']
        dfpmix['Promos$'] = dfpmix['Promotional_Count']*dfpmix['Food_Cost']
        global pmix_total
        pmix_total = dfpmix['Food Cost$'].sum()
        global give_total
        give_total = dfpmix['Promos$'].sum()
        global StopPgm2
        StopPgm2 = ""

#---- READ & TRANSFORM THE DATA ------------------------------------------------------------------` `
st.cache()
def initial_dataread(user_upload):
            

    try:
        dfsales = pd.read_csv(
                    user_upload,
                    )
    except:
        pass

    try:
        sales_calc(dfsales)
    except:
        pass
    
    try:
        pmix_calc(dfsales)
    except:
        pass


    try:
        df = pd.read_excel(
                io=user_upload,
                engine="xlrd",
                sheet_name='InventoryActivityCurrentDayCFA',
                skiprows=12,
            )
    except:
        pass
        try: 
            df = pd.read_excel(
                        io=user_upload,
                        engine="openpyxl",
                        sheet_name='InventoryActivityCurrentDayCFA',
                        skiprows=12,
                    )
        except:
            pass
    
    try:
        #drop all the NaN's and un-named rows
        df.drop(df.columns[df.columns.str.contains('unnamed',case = False)],axis = 1, inplace = True)
        df.dropna(axis=0, inplace=True)

        df.drop(columns='Missing')
        df['Waste'] = df['Cost']*df['= Variance']
        df['Waste'] = df['Waste'].round(decimals=0)
        df['Waste%'] = df['Waste']/(df['- Theoretical Usage']*df['Cost'])
        
        #drop buns and multigrain buns since these are ordered locally
        df = df[~df['Description'].isin(['Buns', 'Bun, Multigrain Brioche'])]

        #re-read the file to find the dates and restaurant name
        try:
            dfdate = pd.read_excel(
                    io=user_upload,
                    engine="xlrd",
                    sheet_name='InventoryActivityCurrentDayCFA',
                    nrows=13
                )
        except ValueError:
            dfdate = pd.read_excel(
                    io=user_upload,
                    engine="openpyxl",
                    sheet_name='InventoryActivityCurrentDayCFA',
                    nrows=13
                ) 
        global restaurant_name
        restaurant_name = dfdate.iat[2,14]

        #Get the dates
        dfdate = dfdate.iat[6,10]
        dfdtsplt = dfdate.split(sep=' ')
        beg_date = dfdtsplt[0]
        end_date = dfdtsplt[4]
        # date_range = beg_date + ' - ' + end_date
        date_list.append(end_date)
        

        # add new columns to the dataframe to identify the dates
        df['Beg. period'] = beg_date
        df['Beg. timestamp'] = pd.to_datetime(dfdtsplt[0])
        df['Beg. timestamp'] = df['Beg. timestamp'].dt.date
        df['Ending period'] = end_date
        df['Month'] = pd.DatetimeIndex(df['Ending period']).month
        month = df['Month Name'] = df['Month'].apply(lambda x: calendar.month_name[x])
        df['Ending timestamp'] = pd.to_datetime(dfdtsplt[4])
        df['Ending timestamp'] = df['Ending timestamp'].dt.date
        df['Period length'] = (df['Ending timestamp']-df['Beg. timestamp']).dt.days
        

        #Add the subcategory column based on internal database
        df2 = pd.read_excel('./assets/cfadatasc.xlsx')

        df = pd.merge(df, 
                        df2, 
                        on ='Description', 
                        how ='inner'
        )

        #Create table with most relevent data
        initialdf = df[['Description', 'Subcategory', 'Waste', 'Waste%','Ending period','Month Name', 'Ending timestamp','Period length', 'Beg. timestamp']]
        #append this data to the df_list
        df_list.append(initialdf)
        global StopPgm3
        StopPgm3 = ""
    except:
        pass





with st.sidebar:
    st.title("Load ActivityReport, PMIX, and Sales Data here")
    user_upload = st.file_uploader(label='Load up to 12 ActivityReports',type=([".xls", ".xlsx", ".csv"]), accept_multiple_files=True)
    if user_upload is not None:
        for uploaded_file in user_upload:
            initial_dataread(uploaded_file)

    # st.title("Step 2: Load PMIX file here")
    # pmix_upload = st.sidebar.file_uploader(label='PMIX for MTD range (ending with most recent count date)',type=([".csv", ".xlsx"]), accept_multiple_files=True)
    # if pmix_upload is not None:
    #     for uploaded_file in pmix_upload:
    #         pmix_dataread(uploaded_file)

    # st.title("Step 3: Load MTD Sales file here")
    # sales_upload = st.sidebar.file_uploader(label='',type=([".csv", ".xlsx"]), accept_multiple_files=True)
    # if sales_upload is not None:
    #         for uploaded_file in sales_upload:
    #             sales_dataread(uploaded_file)

#------------------------------------------------------------------------------------------------------------------------

try:
    concat_df = pd.concat(df_list)
    concat_df = concat_df.reset_index()

    concat_sales_df = pd.concat(salesdf)
    concat_sales_df = concat_sales_df.reset_index()
    sales_total = concat_sales_df["Daily Sales"].sum()
except ValueError:
    pass
#----------------------------------------------------------------------

    #--------------------------------------------------------------------------------------------

    #Defining all the chart data for chart function calls
try:
    dfchicken = concat_df[concat_df['Subcategory'] == 'Chicken'].sort_values(by=['Ending timestamp', 'Description'])
    dfchicken = dfchicken[dfchicken['Period length'] > 5]

    dfbkfst = concat_df[concat_df['Subcategory'] == 'Breakfast'].sort_values(by=['Ending timestamp', 'Description'])
    dfbkfst_stacked = concat_df[concat_df['Subcategory'] == 'Breakfast'].sort_values(by='Ending timestamp')
    dfbkfst = dfbkfst[dfbkfst['Period length'] > 5] #limiting the chart for weekly or monhtly data
    

    dfproduce = concat_df[concat_df['Subcategory'] == 'Produce'].sort_values(by='Ending timestamp')
    dfproduce = dfproduce[dfproduce['Period length'] > 5] #limiting the chart for weekly or monhtly data
    dfproduce_stacked = concat_df[concat_df['Subcategory'] == 'Produce'].sort_values(by='Ending timestamp')

    dfdessert = concat_df[concat_df['Subcategory'] == 'Dessert'].sort_values(by='Ending timestamp')
    dfdessert = dfdessert[dfdessert['Period length'] > 5] #limiting the chart for weekly or monhtly data
    dfdessert_stacked = concat_df[concat_df['Subcategory'] == 'Dessert'].sort_values(by='Ending timestamp')

    dfotherfood = concat_df[concat_df['Subcategory'] == 'Other Food'].sort_values(by='Ending timestamp')
    dfotherfood = dfotherfood[dfotherfood['Period length'] > 5] #limiting the chart for weekly or monhtly data
    dfotherfood_stacked = concat_df[concat_df['Subcategory'] == 'Other Food'].sort_values(by='Ending timestamp')

    dfcoater = concat_df[concat_df['Subcategory'] == 'Coater'].sort_values(by='Ending timestamp')
    dfcoater = dfcoater[dfcoater['Period length'] > 5] #limiting the chart for weekly or monhtly data
    dfcoater_stacked = concat_df[concat_df['Subcategory'] == 'Coater'].sort_values(by='Ending timestamp')

    dfbigop = concat_df[concat_df['Subcategory'].str.contains('Beverage')==False]
    dfbigop = dfbigop[dfbigop['Subcategory'].str.contains('Paper')==False]
    dfbigop = dfbigop[dfbigop['Subcategory'].str.contains('Other Food2')==False]
    dfbigop = dfbigop[dfbigop['Subcategory'].str.contains('Condiments')==False]
    dfbigop = dfbigop[dfbigop['Subcategory'].str.contains('Bread')==False]
    dfbigop = dfbigop[dfbigop['Description'].str.contains('Cone')==False]
    dfbigop = dfbigop[dfbigop['Description'].str.contains('Marinade')==False]
    dfbigop.replace([np.inf, -np.inf], np.nan, inplace=True)
    dfbigop.fillna(value=0, inplace=True)
    dfbigop = dfbigop.drop(dfbigop[dfbigop['Waste%'] < 0].index)


except NameError:
    pass





#----DEFINE BAR CHART COLORS-----------------------------------------------
try:
    colors = {date_list[0]:'#03045e'}
    colors = {date_list[0]:'#03045e',
            date_list[1]:'#023e8a'}
    colors = {date_list[0]:'#03045e',
            date_list[1]:'#023e8a',
            date_list[2]: '#0096c7'}
    colors = {date_list[0]:'#03045e',
            date_list[1]:'#023e8a',
            date_list[2]: '#0096c7',
            date_list[3]: '#48cae4'}
    colors = {date_list[0]:'#03045e',
            date_list[1]:'#023e8a',
            date_list[2]: '#0096c7',
            date_list[3]: '#48cae4',
            date_list[4]: '#ade8f4'}
    colors = {date_list[0]:'#03045e',
            date_list[1]:'#023e8a',
            date_list[2]: '#0096c7',
            date_list[3]: '#48cae4',
            date_list[4]: '#ade8f4',
            date_list[5]: '#7CE9CA'}
    colors = {date_list[0]:'#03045e',
            date_list[1]:'#023e8a',
            date_list[2]: '#0096c7',
            date_list[3]: '#48cae4',
            date_list[4]: '#ade8f4',
            date_list[5]: '#7CE9CA',
            date_list[6]:'#03045e'}
    colors = {date_list[0]:'#03045e',
            date_list[1]:'#023e8a',
            date_list[2]: '#0096c7',
            date_list[3]: '#48cae4',
            date_list[4]: '#ade8f4',
            date_list[5]: '#7CE9CA',
            date_list[6]:'#03045e',
            date_list[7]:'#023e8a'}
    colors = {date_list[0]:'#03045e',
            date_list[1]:'#023e8a',
            date_list[2]: '#0096c7',
            date_list[3]: '#48cae4',
            date_list[4]: '#ade8f4',
            date_list[5]: '#7CE9CA',
            date_list[6]:'#03045e',
            date_list[7]:'#023e8a',
            date_list[8]: '#0096c7'}
    colors = {date_list[0]:'#03045e',
            date_list[1]:'#023e8a',
            date_list[2]: '#0096c7',
            date_list[3]: '#48cae4',
            date_list[4]: '#ade8f4',
            date_list[5]: '#7CE9CA',
            date_list[6]:'#03045e',
            date_list[7]:'#023e8a',
            date_list[8]: '#0096c7',
            date_list[9]: '#48cae4'}
    colors = {date_list[0]:'#03045e',
            date_list[1]:'#023e8a',
            date_list[2]: '#0096c7',
            date_list[3]: '#48cae4',
            date_list[4]: '#ade8f4',
            date_list[5]: '#7CE9CA',
            date_list[6]:'#03045e',
            date_list[7]:'#023e8a',
            date_list[8]: '#0096c7',
            date_list[9]: '#48cae4',
            date_list[10]: '#ade8f4'}
    colors = {date_list[0]:'#03045e',
            date_list[1]:'#023e8a',
            date_list[2]: '#0096c7',
            date_list[3]: '#48cae4',
            date_list[4]: '#ade8f4',
            date_list[5]: '#7CE9CA',
            date_list[6]:'#03045e',
            date_list[7]:'#023e8a',
            date_list[8]: '#0096c7',
            date_list[9]: '#48cae4',
            date_list[10]: '#ade8f4',
            date_list[11]: '#7CE9CA'}
    colors = {date_list[0]:'#03045e',
            date_list[1]:'#023e8a',
            date_list[2]: '#0096c7',
            date_list[3]: '#48cae4',
            date_list[4]: '#ade8f4',
            date_list[5]: '#7CE9CA',
            date_list[6]:'#03045e',
            date_list[7]:'#023e8a',
            date_list[8]: '#0096c7',
            date_list[9]: '#48cae4',
            date_list[10]: '#ade8f4',
            date_list[11]: '#7CE9CA'}
        
except IndexError:
    pass


fig_store = {"Subcategory": [], "Donut": []}


# #----SUBCATEGORY CHART----

try:
    dfsub = concat_df.sort_values(by='Ending timestamp')
    dfsub = dfsub.set_index('Subcategory')
    try:
        dfsub = dfsub.drop(['Chicken'])
        dfsub = dfsub.drop(['Beverage'])
        dfsub = dfsub.drop(['Paper'])
        dfsub = dfsub.drop(['Condiments'])
        dfsub = dfsub.drop(['Other Food2'])
    except:
        pass
    dfsub = dfsub.drop(['Waste%'], axis=1)
    dfsub = dfsub.sort_values(by = ['Ending timestamp','Subcategory'])

    layout = Layout(plot_bgcolor='white')
    config = {'displayModeBar': False}

    subcat_chart=go.Figure(layout=layout)
    for t in dfsub['Ending period'].unique():
        dfp = dfsub[dfsub['Ending period']==t]
        subcat_chart.add_traces(go.Bar(
            x= dfp.index, 
            y = dfp['Waste'], 
            name=t,
            marker_color=colors[t],
            customdata=dfp['Description'],
            ))

    subcat_chart.update_traces(hovertemplate="%{x} <br>Waste: $%{y:.0f} <br> %{customdata}")
    subcat_chart.update_layout(dragmode=False)
    subcat_chart.update_layout(legend=dict(
    orientation="h",
    yanchor="bottom",
    y=1.02,
    xanchor="right",
    x=1))
    subcat_chart.update_layout(width=800, height=600, bargap=0.25)
    subcat_chart.update_yaxes(title_text="Waste $")
    fig_store["Subcategory"] += [subcat_chart]

    # dfsub['Subcategory'] = dfsub['Subcategory'].astype(object)
    dfsubP = dfsub.pivot_table(values='Waste', aggfunc = 'sum', index='Ending timestamp', columns='Subcategory')
    dfsubP2 = dfsub.pivot_table(values='Waste', aggfunc = 'sum', index='Subcategory', columns='Ending timestamp')
except NameError:
    pass
 


#-----DONUT CHART--------------------------------------------------------------------
donut_store = []
def donut_chart(optional=None):
    concatdf = concat_df.drop(concat_df[concat_df['Subcategory'] == 'Paper'].index)
    concatdf = concatdf.drop(concatdf[concatdf['Subcategory'] == 'Beverage'].index)
    concatdf = concatdf.drop(concatdf[concatdf['Subcategory'] == 'Condiments'].index)
    mtdwaste = concatdf['Waste'].sum()
    allowance = sales_total * .0065
    gap = mtdwaste - allowance
    

    colors2 = ['#C71B2A', '#7FC8D9', '#004F71']
    labels = ['Gap','Waste Allowance', 'Net Waste']
    values = [gap, allowance, mtdwaste]

    # Use `hole` to create a donut-like pie chart
    fig = go.Figure(data=[go.Pie(labels=labels, values=values, hole=.6)])
    fig.update_layout(width=650, height=650)
    fig.update_traces(hoverinfo='label+percent', textinfo='label+value', textfont_size=20,
                marker=dict(colors=colors2, line=dict(color='#FFFFFF', width=3)))
    fig.update_traces(showlegend=False)
    donut_store.append(fig)
    if optional == None:
        return st.plotly_chart(fig, use_container_width=True)
    else:
        return None


#-----DEFINE CHICKEN AND COATER CHART STRUCTURES---------------------------------------------------------
colors2 = {'Description':
        ['Chicken, Nuggets',
        'Chicken, Filets',
        'Chicken, Filet Spicy PC',
        'Chicken, Tenders',
        'Chicken, Filets, Grilled',
        'Chicken, Nuggets, Grilled',
        'Chicken, Breakfast Filets',
        'Chicken, Bkfst Filet Spicy'],
        'Color':['#F76AD5','#079FEE','#9958E7','#35D17D','#E6CC91','#868584','#F0ED43','#F78D2A']}
        
chicken_store = []
def chicken_chart(df, optional=None):
    
    dfc = pd.DataFrame(colors2)
    df = dfc.merge(df, on='Description', how='inner')

    config = {'displayModeBar': False}

    layout = Layout(plot_bgcolor='white')

    fig=go.Figure(layout=layout)
    for t in df['Ending period'].unique():
        dfp = df[df['Ending period']==t]
        fig.add_traces(go.Bar(x= dfp['Description'], y = dfp['Waste'], name=t,
                            marker_color=dfp['Color'], hovertemplate="$%{y}%{_xother}"))
    fig.update_layout(width=1450, height=600, bargap=0.15)
    fig.update_traces(showlegend=False)
    fig.update_layout(dragmode=False)
    fig.update_yaxes(title_text="Waste $", gridcolor='#E5E4E3')
    chicken_store.append(fig)
    # st.write('Data with date range of less than 5 days will not be displayed on this chart. This analysis is for a date range of  1 week or longer')
    if optional == None:
        return st.plotly_chart(fig, use_container_width=True, **{'config': config})
    else:
        return None
    
#-------------------------------------------------------------------------------------------------------------



#----Trend Chart--------------------------------------------
trend_store = {"Breakfast": [], "Produce":[], "Dessert": [], "Other Food": []}
breakfast = "Breakfast"
produce = "Produce"
dessert = "Dessert"
otherfood = "Other Food"
chicken = "Chicken"

def trend_chart(df, category, optional=None):

    layout = Layout(plot_bgcolor='white')
    config = {'displayModeBar': False}

    fig=go.Figure(layout=layout)
    for t in df['Ending period'].unique():
        dfp = df[df['Ending period']==t]
        fig.add_traces(go.Bar(x= dfp['Description'], y = dfp['Waste'], name=t,
                            marker_color=colors[t], hovertemplate="%{x} <br>Waste: $%{y}%{_xother}"))
    fig.update_layout(width=1450, height=600, bargap=0.25, xaxis={'categoryorder':'total descending'})
    fig.update_layout(dragmode=False)
    fig.update_layout(legend=dict(
    orientation="h",
    yanchor="bottom",
    y=1.02,
    xanchor="right",
    x=1))
    fig.update_yaxes(title_text="Waste $", gridcolor='#E5E4E3')
    fig.add_annotation(showarrow=False,
    font=dict(size=10), 
    xref='x domain',
    x=0.5,
    yref='y domain',
    y=-0.5
    )
    trend_store[category] += [fig]

    fig2 = go.FigureWidget(fig)

    if optional==None:
        return st.plotly_chart(fig2, use_container_width=True, **{'config': config})
    else:
        return None
    

stack_store = {"Breakfast": [], "Produce":[], "Dessert": [], "Other Food": []}
 

def stacked(df, category, optional=None):

        layout = Layout(plot_bgcolor='white')
        config = {'displayModeBar': False}

        fig = go.Figure(layout=layout)
        for t in df['Ending period'].unique():
            dfp = df[df['Ending period']==t]
            fig.add_traces(go.Bar(x=dfp['Description'], y=dfp['Waste'],name=t, marker_color=colors[t], 
                            texttemplate='$%{y:.0f}'))
        fig.update_layout(barmode='relative', xaxis={'categoryorder':'total descending'})
        fig.update_layout(dragmode=False)
        fig.update_layout(legend=dict(
        orientation="h",
        yanchor="bottom",
        y=1.02,
        xanchor="right",
        x=1))
        fig.update_layout(width=1450, height=600, bargap=0.25)
        fig.update_yaxes(title_text="Waste $", gridcolor='#E5E4E3')
        fig.update_traces(hovertemplate="%{x} <br>Waste: $%{y}%{_xother}")
        stack_store[category] += [fig]
        fig2 = go.FigureWidget(fig)
        
        if optional == None:
            return st.plotly_chart(fig2,use_container_width=True, **{'config': config})
        else:
            return None

table_list = {"Chicken": [],"Breakfast": [], "Produce": [], "Dessert": [], "Other Food": []}

def data_tables_pptx(df, category):
    pivot = df.pivot_table(
        index='Description', 
        columns='Ending timestamp', 
        values='Waste', aggfunc='mean')
    df2 = pd.DataFrame(pivot)
    df2.fillna(value=0, inplace=True)
    df2 = df2.astype(float).applymap('${:.1f}'.format)
    table1 = df2.style.highlight_max(color='#CC0003', axis=1)
    table_list[category] += [table1]

    
    pivot2 = df.pivot_table(index='Description', columns='Ending timestamp', values='Waste%')
    df3 = pd.DataFrame(pivot2)
    df3.replace([np.inf, -np.inf], np.nan, inplace=True)
    df3.fillna(value=0, inplace=True)
    df3[df3 < 0] = 0
    df3 = df3[df3.select_dtypes(include=['number']).columns] #*= 100
    df3 = df3.astype(float).applymap('{:.1%}'.format)
    table2 = df3.style.background_gradient(axis=1, cmap='PuBu')
    table_list[category] += [table2]

def styling(v, cmap=''):
        try:
            return cmap if v > 0 else None
        except:
            pass

def data_tables(df):
    col1, col2 = st.columns(2)
    with col1:
        st.title('DataTable $ Wasted',)
        pivot = df.pivot_table(
            index='Description', 
            columns='Ending timestamp', 
            values='Waste', aggfunc='mean')
        df2 = pd.DataFrame(pivot)
        df2.fillna(value=0, inplace=True)
        # table1 = df2.astype(float).applymap('${:.0f}'.format)
        df2_to_dlr = {}
        for i in df2:
            df2_to_dlr[i] = '${:.0f}'.format
        st.table(df2.style.hide().applymap(styling, cmap='color:Blue;').background_gradient(axis=1, cmap='Reds').format(df2_to_dlr))
        # st.table(table1)


    with col2:
        st.title('% of Used Product Wasted')
        pivot2 = df.pivot_table(index='Description', columns='Ending timestamp', values='Waste%')
        df3 = pd.DataFrame(pivot2)
        df3.replace([np.inf, -np.inf], np.nan, inplace=True)
        df3.fillna(value=0, inplace=True)
        df3[df3 < 0] = 0
        df3 = df3[df3.select_dtypes(include=['number']).columns] #*= 100
        df3_to_pct = {}
        for i in df3:
            df3_to_pct[i] = '{:.1%}'.format
        st.table(df3.style.hide().applymap(styling, cmap='color:Blue;').background_gradient(axis=1, cmap='PuBu').format(df3_to_pct))
  
def big_op(df):
        df.drop(df[df['Waste%'] >= 1].index, inplace = True)
        df['Waste%%'] = df['Waste%']*100
        fig = px.scatter(df, 
            x='Description', 
            y="Waste%%",
            color="Subcategory", 
            width=1400, 
            height=700, 
            size='Waste%',
            hover_name="Description", 
            size_max=60,
            custom_data=['Ending period'])
        fig.update_traces(mode='markers', hovertemplate="%{x} <br>Waste Percentage: %{y:.0f}% <br> %{customdata}")
        fig.update_yaxes(gridcolor='#E5E4E3')
        fig.update_layout(plot_bgcolor='white')
        #fig.update_xaxes(rangeslider_visible=True)
        st.plotly_chart(fig, use_container_width=True)
        st.markdown('Waste percentages are calculated by dividing total dollars in waste by total dollars spent on used inventory (theoretical * cost/case). Think of the percentage of waste number (ex. 10%) as the average percentage of product in a given case that was wasted (ex. on average, 10% of a case of chicken was wasted during the measured period). Ideally, we would want to see less than 20% for most items.')

    #------------------------



if selected == "DataViz":
    st.title("DataViz")
    try:
        StopPgm3
    except NameError:
        st.subheader("Upload ActivityReport data to continue")
        st.stop()
    p_selection = option_menu(
        menu_title=None,
        options=["Subcategory", "Chicken", "Breakfast", "Produce", "Dessert", "Other Food", "Percentages"],
        icons=["bar-chart-steps", "bar-chart-fill", "bar-chart-fill", "bar-chart-fill", "bar-chart-fill", "bar-chart-fill", "bar-chart-fill"],
        orientation="horizontal"
    )
    if p_selection == "Subcategory":
        col1, col2 = st.columns(2)
        with col1:
            st.markdown("<h1 style='text-align: center; color: #073D8D;'>Subcategory Chart</h1>", unsafe_allow_html=True)
            st.plotly_chart(subcat_chart, use_container_width=True, **{'config':config})
            st.markdown("<h1 style='text-align: center; color: #073D8D;'>DataTable $'s Wasted</h1>", unsafe_allow_html=True)
            sub_table = dfsubP2
            df_to_dl = {}
            for i in sub_table:
                df_to_dl[i] = '${:.0f}'.format
            st.table(sub_table.style.hide().applymap(styling, cmap='color:Blue;').background_gradient(axis=1, cmap='PuBu').format(df_to_dl))
        with col2:
            try:
                st.markdown("<h1 style='text-align: center; color: #073D8D;'>Net Food Cost Gap</h1>", unsafe_allow_html=True)
                donut_chart()
            except NameError:
                st.markdown("<h1 style='text-align: center; color: #00000;'>No Sales Data Uploaded</h1>", unsafe_allow_html=True)
            st.markdown("<h1 style='text-align: center; color: #02305F;'>Total waste by period</h1>", unsafe_allow_html=True)
            df = dfbigop.groupby(["Ending timestamp", "Ending period"])["Waste"].sum().reset_index()
            df_latest = df['Ending timestamp'].max()
            st.table(df)
            
            
            
            
    if p_selection == "Chicken":
        st.markdown("<h1 style='text-align: left; color: #C60C31;'>Chicken Waste Trend</h1>", unsafe_allow_html=True)
        chicken_chart(dfchicken)
        data_tables(dfchicken)
    if p_selection == "Breakfast":
        st.markdown("<h1 style='text-align: left; color: #C60C31;'>Breakfast Waste Trend</h1>", unsafe_allow_html=True)
        trend_chart(dfbkfst, breakfast)
        st.markdown("<h1 style='text-align: left; color: #C60C31;'>Breakfast Waste Stacked</h1>", unsafe_allow_html=True)
        stacked(dfbkfst, breakfast)
        data_tables(dfbkfst)
    else:
        ""
    if p_selection == "Produce":
        st.markdown("<h1 style='text-align: left; color: #C60C31;'>Produce Waste Trend</h1>", unsafe_allow_html=True)
        trend_chart(dfproduce, produce)
        st.markdown("<h1 style='text-align: left; color: #C60C31;'>Produce Waste Stacked</h1>", unsafe_allow_html=True)
        stacked(dfproduce_stacked, produce)
        data_tables(dfproduce)
    else:
        ""
    if p_selection == "Dessert":
        st.markdown("<h1 style='text-align: left; color: #C60C31;'>Dessert Waste Trend</h1>", unsafe_allow_html=True)
        trend_chart(dfdessert, dessert)
        st.markdown("<h1 style='text-align: left; color: #C60C31;'>Dessert Waste Stacked</h1>", unsafe_allow_html=True)
        stacked(dfdessert_stacked, dessert)
        data_tables(dfdessert)
    else:
        ""
    if p_selection == "Other Food":
        st.markdown("<h1 style='text-align: left; color: #C60C31;'>Other Food Waste Trend</h1>", unsafe_allow_html=True)
        trend_chart(dfotherfood, otherfood)
        st.markdown("<h1 style='text-align: left; color: #C60C31;'>Other Food Waste Stacked</h1>", unsafe_allow_html=True)
        stacked(dfotherfood_stacked, otherfood)
        data_tables(dfotherfood)
    else:
        ""
    if p_selection == "Percentages":
        st.markdown("<h1 style='text-align: left; color: #C60C31;'>Waste Percentages</h1>", unsafe_allow_html=True)
        big_op(dfbigop)
    else:
        ""






if selected == "Food Cost %":

    st.title("Calculate Your Food Cost Percentage")
    try:
        StopPgm
        StopPgm2
    except NameError:
        st.subheader("Upload ActivityReports, PMIX, and Sales data to continue")
        st.stop()


    st.markdown(":bar_chart:" + restaurant_name)
    st.subheader("Select FSU / Mall")
    mode = st.selectbox(label='', options=('FSU', 'Mall'))
    st.subheader("Select the month you are calulating your food cost for ")
    option = st.selectbox(
            "note: if your end of month count was submitted on the 1st of the following month, you will need to mannually change the submission date in the excel file to the last day of the month",
            ('January', 'February', 'March', 'April', 'May', 'June', 'July', 'August', 'September', 'October', 'November', 'December'))

    st.subheader(f'You selected: {option}')

    
    concatdf = concat_df.drop(concat_df[concat_df['Subcategory'] == 'Paper'].index)
    concatdf = concatdf.drop(concatdf[concatdf['Subcategory'] == 'Beverage'].index)
    concatdf = concatdf.drop(concatdf[concatdf['Subcategory'] == 'Condiments'].index)
    concatdf['Total'] = concatdf.loc[concatdf['Month Name'] == option,['Waste']].sum(axis=1)
    concatdf['Total'].fillna(0, inplace=True)
    mtdwaste = concatdf['Total'].sum()
    concat_sales_df['Total2'] = concat_sales_df.loc[concat_sales_df['Month Name'] == option,['Daily Sales']].sum(axis=1)
    concat_sales_df['Total2'].fillna(0, inplace=True)
    month_sales = concat_sales_df['Total2'].sum()

    
    col1, col2, col3, col4, col5 = st.columns(5)
    with col1:
                
        
        with st.form(key='my_form'):
            pmsales = st.number_input('Prior month sales')
            pmrefill = st.number_input('Prior month cost of refills (on the food cost report)')
            discounts = st.number_input('MTD Discounts (run the current month marketing and giveaways report in inform)')
            rebates = st.number_input('Beverage rebates ~200-300 per week')
            credit = st.number_input('Reimbursments receieved this month from QIC submissions')
            submitted = st.form_submit_button(label='Submit')
            
            
            st.cache()
            def calc():
                
                x = "{value:,.0f}"
                y = "{value:.2f}"


                if mode == "Mall":
                    waste_allowance = month_sales*.0105
                else:
                    waste_allowance = month_sales*.0065
            

                #target food cost calculation
                condiments = month_sales*.0141
                # waste_allowance = month_sales*.0065
                refills = pmrefill*month_sales/pmsales
                global target
                target = pmix_total + waste_allowance + condiments + refills - discounts
                targetP = target / month_sales*100
        
                
                
                    


                #Estimated food cost
                total_waste = concat_df['Waste'].sum()
                beverage = concat_df[concat_df['Subcategory'] == 'Beverage']['Waste'].sum()
                global food_cost
                food_cost = pmix_total + mtdwaste + condiments - discounts - credit - rebates + refills
                estfcP = target / month_sales*100
                

                
            
                with col2:
                    st.subheader("PMIX Total")
                    st.subheader("MTD Waste")
                    st.subheader("Condiments")
                    st.subheader("Discounts")
                    st.subheader("Rebates")
                    st.subheader("QIC Credits")
                    st.subheader("Refills")
                    st.subheader("Food Cost $")
                    st.subheader("Sales Total")
                    st.subheader("Est. FC %")
                with col3:
                    st.subheader("+ $" + x.format(value=(pmix_total)))
                    st.subheader("+ $" + x.format(value=(mtdwaste)))
                    st.subheader("+ $" + x.format(value=(condiments)))
                    st.subheader("- $" + x.format(value=(discounts)))
                    st.subheader("- $" + x.format(value=(rebates)))
                    st.subheader("- $" + x.format(value=(credit)))
                    st.subheader("+ $" + x.format(value=(refills)))
                    st.subheader("$" + x.format(value=(food_cost)))
                    st.subheader("$" + x.format(value=(month_sales)))
                    st.subheader(y.format(value=(food_cost / month_sales * 100)) + '%')
                with col4:
                    st.subheader("PMIX Total")
                    st.subheader("Waste Allowance")
                    st.subheader("Condiments")
                    st.subheader("Discounts")
                    st.markdown('---')
                    st.markdown('---')
                    st.subheader("Refills")
                    st.subheader("Target FC $")
                    st.subheader("Sales Total")
                    st.subheader("Target FC %")
                with col5:
                    st.subheader("+ $" + x.format(value=(pmix_total)))
                    st.subheader("+ $" + x.format(value=(waste_allowance)))
                    st.subheader("+ $" + x.format(value=(condiments)))
                    st.subheader("- $" + x.format(value=(discounts)))
                    st.markdown('---')
                    st.markdown('---')
                    st.subheader("+ $" + x.format(value=(refills)))
                    st.subheader("$" + x.format(value=(target)))
                    st.subheader("$" + x.format(value=(month_sales)))
                    st.subheader(y.format(value=(target / month_sales * 100)) + '%')

    
    
                # st.title("Download Your Food Cost Report")
                prs=Presentation('./assets/testpresentation.pptx')
                
                slide0 = prs.slides[0]
                title = slide0.shapes.title
                concat_df_date = concat_df
                most_recent_date = concat_df_date.sort_values(by='Ending timestamp', inplace=True)
                most_recent_date = concat_df_date['Ending period'].iloc[-1]
                title.text = "As of " + most_recent_date
                title.text_frame.paragraphs[0].font.color.rgb = RGBColor(198, 12, 49)
                text1 = slide0.shapes.placeholders[10]
                text1.text = restaurant_name
                blank_slide_layout = prs.slide_layouts[0]


                est = "% s" % (y.format(value=(target / month_sales * 100)) + '%')
                targ = "% s" % (y.format(value=(food_cost / month_sales * 100)) + '%')

                fc_slide = prs.slides[1]
                text1 = fc_slide.shapes.placeholders[15]
                text1.text = est
                text1 = fc_slide.shapes.placeholders[16]
                text1.text = targ
                text1 = fc_slide.shapes.placeholders[17]
                text1.text = most_recent_date
                text1 = fc_slide.shapes.placeholders[18]
                text1.text = most_recent_date





                #slide titles
                chicken_waste = "Chicken Waste"
                bkfst_waste = "Breakfast Waste"
                produce_waste = "Produce Waste"
                dessert_waste = "Dessert Waste"
                otherfood_waste = "Other Food Waste"

                ph = "write chart to dictionairy"

                try:
                    donut_chart(1)
                except:
                    pass

                chicken_chart(dfchicken, ph)
                data_tables_pptx(dfchicken, chicken)

                trend_chart(dfbkfst, breakfast, ph)
                stacked(dfbkfst_stacked, breakfast, ph)
                data_tables_pptx(dfbkfst, breakfast)

                trend_chart(dfproduce, produce, ph)
                stacked(dfproduce_stacked, produce, ph)
                data_tables_pptx(dfproduce, produce)

                trend_chart(dfdessert, dessert, ph)
                stacked(dfdessert_stacked, dessert, ph)
                data_tables_pptx(dfdessert, dessert)

                trend_chart(dfotherfood, otherfood, ph)
                stacked(dfotherfood_stacked, otherfood, ph)
                data_tables_pptx(dfotherfood, otherfood)

                #label charts using dictionaries
                subcategory = fig_store['Subcategory'][0]
                breakfast_tchart = trend_store['Breakfast'][0]
                breakfast_stchart = stack_store['Breakfast'][0]
                try:
                    chicken_table1 = table_list['Chicken'][0]
                    chicken_table2 = table_list['Chicken'][1]
                    breakfast_table1 = table_list['Breakfast'][0]
                    breakfast_table2 = table_list['Breakfast'][1]
                    produce_table1 = table_list['Produce'][0]
                    produce_table2 = table_list['Produce'][1]
                    dessert_table1 = table_list['Dessert'][0]
                    dessert_table2 = table_list['Dessert'][1]
                    otherfood_table1 = table_list['Other Food'][0]
                    otherfood_table2 = table_list['Other Food'][1]
                except:
                    pass
                chicken_chart_image = chicken_store[0]
                produce_tchart = trend_store['Produce'][0]
                produce_stchart = stack_store['Produce'][0]
                dessert_tchart = trend_store['Dessert'][0]
                dessert_stchart = stack_store['Dessert'][0]
                otherfood_tchart = trend_store['Other Food'][0]
                otherfood_stchart = stack_store['Other Food'][0]


                def summary_chart_slide(chart1):
                    with NamedTemporaryFile(delete=False, suffix=".png") as tmpfile:
                            chart1.write_image(tmpfile.name)
                            path = tmpfile.name
                    try:
                        with NamedTemporaryFile(delete=False, suffix=".png") as tmpfile2:
                                donut_store[0].write_image(tmpfile2.name)
                                path2 = tmpfile2.name
                    except IndexError:
                        pass

                    slide = prs.slides.add_slide(blank_slide_layout)
                    title = slide.shapes.title
                    title.text = "Waste Summary"
                    left=Inches(.25)
                    top=Inches(1.5)
                    img=slide.shapes.add_picture(path,left, top, width=Inches(6.15))

                    try:
                        left=Inches(6.3)
                        top=Inches(2.2)
                        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(198, 12, 49)
                        img=slide.shapes.add_picture(path2,left, top, width=Inches(3.6))
                    except:
                        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(198, 12, 49)
                        txBox = slide.shapes.add_textbox(left=Inches(6.3), top=Inches(2.2),
                                            width=Inches(1),height=Inches(1))
                        tf = txBox.text_frame
                        tf.text = "No Sales Data Uploaded"

                    tmpfile.close()
                    os.unlink(tmpfile.name)
                    tmpfile2.close()
                    os.unlink(tmpfile2.name)


                def single_chart_slide(chart1, table1, table2, slide_title):
                    with NamedTemporaryFile(delete=False, suffix=".png") as tmpfile:
                            pio.write_image(chart1, tmpfile.name, engine="kaleido")
                            path = tmpfile.name
                    with NamedTemporaryFile(delete=False, suffix=".png") as tmpfile2:
                            dfi.export(table1, tmpfile2.name, table_conversion='matplotlib')
                            path3 = tmpfile2.name
                    with NamedTemporaryFile(delete=False, suffix=".png") as tmpfile3:
                            dfi.export(table2, tmpfile3.name, table_conversion='matplotlib')
                            path4 = tmpfile3.name
                    
                    slide = prs.slides.add_slide(blank_slide_layout)
                    title = slide.shapes.title
                    title.text = slide_title
                    left=Inches(1)
                    top=Inches(1.35)
                    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(198, 12, 49)
                    img=slide.shapes.add_picture(path,left, top, width=Inches(8))
                    left=Inches(.5)
                    top=Inches(4.7)
                    img=slide.shapes.add_picture(path3,left, top, width=Inches(4))
                    left=Inches(5.5)
                    top=Inches(4.7)
                    img=slide.shapes.add_picture(path4,left, top, width=Inches(4))
                    
                    tmpfile.close()
                    os.unlink(tmpfile.name)
                    tmpfile2.close()
                    os.unlink(tmpfile2.name)
                    tmpfile3.close()
                    os.unlink(tmpfile3.name)

                def multi_chart_slide(chart1,chart2,table1, table2, slide_title, wdth):
                    with NamedTemporaryFile(delete=False, suffix=".png") as tmpfile1:
                            pio.write_image(chart1, tmpfile1.name, engine="kaleido")
                            path1 = tmpfile1.name
                    with NamedTemporaryFile(delete=False, suffix=".png") as tmpfile2:
                            pio.write_image(chart2, tmpfile2.name, engine="kaleido")
                            path2 = tmpfile2.name
                    with NamedTemporaryFile(delete=False, suffix=".png") as tmpfile3:
                            dfi.export(table1, tmpfile3.name, table_conversion='matplotlib')
                            path3 = tmpfile3.name
                    with NamedTemporaryFile(delete=False, suffix=".png") as tmpfile4:
                            dfi.export(table2, tmpfile4.name, table_conversion='matplotlib')
                            path4 = tmpfile4.name

                    slide = prs.slides.add_slide(blank_slide_layout)
                    title = slide.shapes.title
                    title.text = slide_title
                    left=Inches(.25)
                    top=Inches(1.25)
                    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(198, 12, 49)
                    img=slide.shapes.add_picture(path1,left, top, width=Inches(6))
                    left=Inches(.25)
                    top=Inches(4)
                    img=slide.shapes.add_picture(path2,left, top, width=Inches(6))
                    left=Inches(6.45)
                    top=Inches(1.45)
                    img=slide.shapes.add_picture(path3,left, top, width=Inches(wdth))
                    left=Inches(6.45)
                    top=Inches(4.2)
                    img=slide.shapes.add_picture(path4,left, top, width=Inches(wdth))

                    tmpfile1.close()
                    os.unlink(tmpfile1.name)
                    tmpfile2.close()
                    os.unlink(tmpfile2.name)
                    tmpfile3.close()
                    os.unlink(tmpfile3.name)
                    tmpfile4.close()
                    os.unlink(tmpfile4.name)



                summary_chart_slide(subcategory)
                single_chart_slide(chicken_chart_image, chicken_table1, chicken_table2, chicken_waste)
                multi_chart_slide(breakfast_tchart, breakfast_stchart, breakfast_table1, breakfast_table2, bkfst_waste, 3.3)
                multi_chart_slide(produce_tchart, produce_stchart, produce_table1, produce_table2, produce_waste, 3.3)
                multi_chart_slide(dessert_tchart, dessert_stchart, dessert_table1, dessert_table2, dessert_waste, 3.3)
                multi_chart_slide(otherfood_tchart, otherfood_stchart, otherfood_table1, otherfood_table2, otherfood_waste, 3.3)
                
                with col1:
                    with NamedTemporaryFile(delete=False, suffix=".pptx") as tmpfile:
                        prs.save(tmpfile.name) 
                        pptx = tmpfile.name
                        def get_binary_file_downloader_html(bin_file, file_label='File'):
                            with open(bin_file, 'rb') as f:
                                btn = st.download_button(
                                label=f"Download PowerPoint Report for {option}",
                                data=f.read(),
                                file_name="PowerPoint Report.pptx" # Any file name
                            )
                        get_binary_file_downloader_html(pptx, 'Powerpoint Report')

                        tmpfile.close()
                        os.unlink(tmpfile.name)

        
    if submitted:
        st.button(label=f"Calculate Food Cost Percentage for {option}", on_click=calc)
    if not submitted:
        st.stop()
    # st.button(label=f"Download Report for {option}", on_click=pptx)
    # st.dataframe(concat_df)
    # st.dataframe(concat_sales_df)